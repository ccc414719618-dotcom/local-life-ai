"""
Quinn Finance Report Generator - Professional PPTX
生成图文并茂的专业金融分析报告
"""

import os
import io
import sys
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
import numpy as np
from datetime import datetime
from PIL import Image

# 中文字体
plt.rcParams['font.sans-serif'] = ['PingFang SC', 'Hiragino Sans GB', 'Microsoft YaHei', 'Arial Unicode MS', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

WORKSPACE = '/Volumes/1TB/openclaw/jinrong-bot/finance'
OUTPUT_DIR = f'{WORKSPACE}/reports'
os.makedirs(OUTPUT_DIR, exist_ok=True)

from chart_generator import ChartGenerator
from data_fetcher import DataFetcher
from report_generator import TechnicalAnalyzer


def fig_to_image(fig, dpi=120):
    """matplotlib figure → PIL Image"""
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
               facecolor=fig.get_facecolor(), edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return Image.open(buf).convert('RGB')


def _make_table_image(headers, rows, col_widths=None, title='',
                      header_color='#1C1C1C', row_colors=('#FFFFFF', '#F5F5F5'),
                      fontsize=10, cell_height=28, width=800):
    """生成表格图片"""
    n_cols = len(headers)
    n_rows = len(rows)
    if col_widths is None:
        col_widths = [width // n_cols] * n_cols
    total_width = sum(col_widths)
    
    fig_h = cell_height * (n_rows + 2) + (30 if title else 0)
    fig, ax = plt.subplots(figsize=(total_width / 100, fig_h / 100), dpi=100)
    fig.patch.set_facecolor('white')
    ax.set_facecolor('white')
    
    y = 1.0
    # 标题
    if title:
        ax.text(total_width / 2, y, title, ha='center', va='top',
               fontsize=fontsize + 2, fontweight='bold', color='#1C1C1C')
        y -= 0.05
    
    # 表头
    x_positions = [sum(col_widths[:i]) + col_widths[i]/2 for i in range(n_cols)]
    for i, (h, xp) in enumerate(zip(headers, x_positions)):
        rect = plt.Rectangle((xp - col_widths[i]/2, y - cell_height/100), col_widths[i], cell_height/100,
                             facecolor='#2D2D2D', edgecolor='white', linewidth=0.5)
        ax.add_patch(rect)
        ax.text(xp, y - cell_height/200, h, ha='center', va='center',
               fontsize=fontsize, fontweight='bold', color='white')
    y -= cell_height / 100
    
    # 数据行
    for ri, row in enumerate(rows):
        bg_color = row_colors[ri % len(row_colors)]
        for ci, (cell, xp) in enumerate(zip(row, x_positions)):
            rect = plt.Rectangle((xp - col_widths[ci]/2, y - cell_height/100), col_widths[ci], cell_height/100,
                                 facecolor=bg_color, edgecolor='#EEEEEE', linewidth=0.3)
            ax.add_patch(rect)
            ax.text(xp, y - cell_height/200, str(cell), ha='center', va='center',
                   fontsize=fontsize, color='#333333')
        y -= cell_height / 100
    
    ax.set_xlim(0, total_width)
    ax.set_ylim(y - 0.05, 1.0)
    ax.axis('off')
    plt.tight_layout(pad=0)
    return fig_to_image(fig, dpi=100)


def _fmt_price(price):
    if price is None: return 'N/A'
    try:
        if abs(price) >= 1000:
            return f'${price:,.0f}'
        else:
            return f'${price:,.2f}'
    except (TypeError, ValueError):
        return 'N/A'


def _fmt_val(val, spec=',.2f', fallback='N/A'):
    """安全格式化：None/非数值返回fallback"""
    if val is None: return fallback
    try:
        return f'{float(val):{spec}}'
    except (TypeError, ValueError):
        return fallback


def _fmt_change(change):
    return f'{change:+.2f}%'


class QuinnReportPPTX:
    """生成专业PPT报告"""
    
    def __init__(self, fetcher=None):
        # 复用外部传入的 DataFetcher 实例（利用其_yfetch缓存避免重复下载）
        self.fetcher = fetcher if fetcher else DataFetcher()
        self.analyzer = TechnicalAnalyzer()
        self.today = datetime.now().strftime('%Y-%m-%d')
        self.charts_dir = f'{WORKSPACE}/charts'
        os.makedirs(self.charts_dir, exist_ok=True)
    
    def _get_data(self):
        """获取所有数据"""
        crypto = self.fetcher.get_crypto_prices(['bitcoin', 'ethereum', 'solana'])
        ohlc = self.fetcher.get_crypto_ohlc('bitcoin', 90)
        closes = [k[4] for k in ohlc] if ohlc else []
        highs = [k[2] for k in ohlc] if ohlc else []
        lows = [k[3] for k in ohlc] if ohlc else []
        
        ma5 = self.analyzer.calculate_ma(closes, 5)
        ma10 = self.analyzer.calculate_ma(closes, 10)
        ma20 = self.analyzer.calculate_ma(closes, 20)
        rsi = self.analyzer.calculate_rsi(closes)
        macd, signal_line, histogram = self.analyzer.calculate_macd(closes)
        bb_upper, bb_middle, bb_lower = self.analyzer.calculate_bollinger(closes)
        
        metals = self.fetcher.get_precious_metals()
        indices = self.fetcher.get_market_indices()
        fear_greed = self.fetcher.get_fear_greed_index()
        news = self.fetcher.get_news(['BTC-USD', '^GSPC', 'GC=F'], limit=5)
        
        current = closes[-1] if closes else 0
        prev = closes[-2] if len(closes) > 1 else current
        change_pct = ((current - prev) / prev * 100) if prev else 0
        
        # 趋势
        if ma20 and ma5:
            trend = '上升趋势' if current > ma20 and ma20 > ma5 else ('下降趋势' if current < ma20 else '震荡整理')
        elif ma20:
            trend = '上升趋势' if current > ma20 else ('下降趋势' if current < ma20 else '震荡整理')
        else:
            trend = '震荡整理'
        
        # RSI 状态
        if rsi:
            if rsi > 70: rsi_status = '超买，警惕回调'
            elif rsi < 30: rsi_status = '超卖，关注反弹'
            elif rsi > 60: rsi_status = '偏强'
            elif rsi < 40: rsi_status = '偏弱'
            else: rsi_status = '中性'
        else:
            rsi_status = '数据不足'
        
        # MACD 状态
        if macd and histogram is not None:
            macd_status = '多方动能增强' if histogram > 0 else '空方动能增强'
        else:
            macd_status = '数据不足'
        
        # 综合信号
        signals = []
        if current > ma20: signals.append(True)
        else: signals.append(False)
        if rsi and rsi < 30: signals.append(True)
        elif rsi and rsi > 70: signals.append(False)
        else: signals.append(None)
        if histogram is not None and histogram > 0: signals.append(True)
        elif histogram is not None: signals.append(False)
        else: signals.append(None)
        bull = sum(1 for s in signals if s is True)
        bear = sum(1 for s in signals if s is False)
        if bull > bear and bull >= 2: overall = '偏多'
        elif bear > bull and bear >= 2: overall = '偏空'
        else: overall = '中性'
        
        # 支撑阻力
        recent_high = max(highs[-7:]) if len(highs) >= 7 else max(highs) if highs else 0
        recent_low = min(lows[-7:]) if len(lows) >= 7 else min(lows) if lows else 0
        pivot = (highs[-1] + lows[-1] + closes[-1]) / 3 if highs and lows and closes else 0
        
        return {
            'crypto': crypto,
            'metals': metals,
            'indices': indices,
            'fear_greed': fear_greed,
            'news': news,
            'closes': closes,
            'highs': highs,
            'lows': lows,
            'ma5': ma5, 'ma10': ma10, 'ma20': ma20,
            'rsi': rsi, 'macd': macd, 'signal_line': signal_line,
            'histogram': histogram,
            'bb_upper': bb_upper, 'bb_middle': bb_middle, 'bb_lower': bb_lower,
            'current': current, 'prev': prev, 'change_pct': change_pct,
            'trend': trend, 'rsi_status': rsi_status, 'macd_status': macd_status,
            'overall': overall, 'bull': bull, 'bear': bear,
            'recent_high': recent_high, 'recent_low': recent_low, 'pivot': pivot,
        }
    
    def _build_cover_slide(self, prs, data):
        """封面"""
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.oxml.ns import qn
        from pptx.oxml import parse_xml
        
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        
        # 背景
        background = slide.shapes.add_shape(
            1, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(0x0D, 0x11, 0x17)
        background.line.fill.background()
        
        # 顶部装饰条
        top_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.1))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = RGBColor(0xF7, 0x93, 0x1A)  # 橙色
        top_bar.line.fill.background()
        
        # 主标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.67), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = 'QUINN 专业金融分析报告'
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.67), Inches(0.6))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        btc_price = data['crypto'].get('bitcoin', {}).get('usd', 0)
        btc_change = data['crypto'].get('bitcoin', {}).get('change_24h', 0)
        p2.text = f'BTC ${btc_price:,.0f}  ({btc_change:+.2f}%)  ·  {self.today}'
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0xF7, 0x93, 0x1A)
        p2.alignment = PP_ALIGN.CENTER
        
        # 底部信息
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.67), Inches(0.5))
        tf3 = footer_box.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = 'Quinn 投资研究员  ·  加密货币 & 贵金属 & 全球市场'
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        p3.alignment = PP_ALIGN.CENTER
    
    def _add_slide_header(self, slide, prs, title, subtitle=''):
        """通用幻灯片头部"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        # 顶部色条
        bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(0xF7, 0x93, 0x1A)
        bar.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1C, 0x1C, 0x1C)
        
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.6), Inches(12), Inches(0.35))
            tf2 = sub_box.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(11)
            p2.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    
    def _add_image_to_slide(self, slide, img_path, left, top, width, height):
        """添加图片到幻灯片"""
        from pptx.util import Inches
        from PIL import Image as PILImage
        
        if not os.path.exists(img_path):
            return
        img = PILImage.open(img_path)
        w, h = img.size
        aspect = w / h
        if width / height > aspect:
            width = height * aspect
        else:
            height = width / aspect
        slide.shapes.add_picture(img_path, left, top, width, height)
    
    def _add_text_box(self, slide, text, left, top, width, height,
                      font_size=11, bold=False, color=(0x33, 0x33, 0x33),
                      align='left', bg_color=None):
        """添加文本框"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
        if align == 'center': p.alignment = PP_ALIGN.CENTER
        elif align == 'right': p.alignment = PP_ALIGN.RIGHT
        
        if bg_color:
            fill = txBox.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*bg_color)
    
    def _build_market_overview_slide(self, prs, data):
        """市场总览幻灯片"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.oxml.ns import qn
        import copy
        
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._add_slide_header(slide, prs, '📊 市场总览', f'报告日期: {self.today}')
        
        # 左：加密货币表格
        headers = ['品种', '价格 (USD)', '24h涨跌']
        rows = []
        for sym, info in data['crypto'].items():
            change = info['change_24h']
            emoji = '🔴' if change < 0 else '🟢'
            price_str = _fmt_price(info['usd'])
            rows.append([sym.capitalize(), price_str, f'{emoji} {change:+.2f}%'])
        
        fig_table = _make_table_image(headers, rows, title='加密货币',
                                      col_widths=[180, 200, 160],
                                      header_color='#1C1C1C')
        img_path = f'{self.charts_dir}/crypto_table.png'
        fig_table.save(img_path)
        self._add_image_to_slide(slide, img_path, Inches(0.3), Inches(1.0), Inches(5.0), Inches(1.5))
        
        # 中：贵金属表格
        headers2 = ['品种', '价格 (USD/盎司)', '24h涨跌']
        rows2 = []
        for name, info in data['metals'].items():
            change = info['change']
            emoji = '🔴' if change < 0 else '🟢'
            rows2.append([name, _fmt_price(info['price']), f'{emoji} {change:+.2f}%'])
        
        fig_m = _make_table_image(headers2, rows2, title='贵金属', col_widths=[200, 220, 150])
        img_path2 = f'{self.charts_dir}/metals_table.png'
        fig_m.save(img_path2)
        self._add_image_to_slide(slide, img_path2, Inches(5.5), Inches(1.0), Inches(5.0), Inches(1.2))
        
        # 右：恐慌贪婪
        fg = data['fear_greed']
        if fg:
            if fg >= 75: fg_label = '极度贪婪', (0xFF, 0x44, 0x44)
            elif fg >= 50: fg_label = '贪婪', (0xFF, 0xCC, 0x00)
            elif fg >= 25: fg_label = '恐惧', (0x44, 0x88, 0xFF)
            else: fg_label = '极度恐惧', (0x44, 0x44, 0x44)
            
            fg_box = slide.shapes.add_shape(1, Inches(5.5), Inches(2.3), Inches(5.0), Inches(1.2))
            fg_box.fill.solid()
            fg_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            fg_box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        
        # 底部：市场涨跌柱状图
        chart_path = ChartGenerator.plot_market_overview(
            data['crypto'], data['metals'], data['indices'])
        if chart_path:
            self._add_image_to_slide(slide, chart_path, Inches(0.3), Inches(3.5), Inches(12.67), Inches(3.8))
    
    def _build_global_indices_slide(self, prs, data):
        """全球指数幻灯片"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import copy
        
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._add_slide_header(slide, prs, '🌍 全球主要指数', '美国 · 欧洲 · 亚太')
        
        indices = data['indices']
        if not indices:
            self._add_text_box(slide, '暂无指数数据', Inches(1), Inches(2), Inches(10), Inches(1),
                             font_size=14, color=(0x88, 0x88, 0x88), align='center')
            return
        
        # 分区域展示
        regions = {
            '美国': ['S&P 500', 'NASDAQ', 'Dow Jones'],
            '欧洲': ['FTSE 100', 'DAX'],
            '亚太': ['Nikkei 225', 'Hang Seng', 'SSE 上证'],
        }
        
        y_offset = 1.1
        for region, names in regions.items():
            # 区域标题
            self._add_text_box(slide, f'  {region}', Inches(0.3), Inches(y_offset),
                             Inches(3), Inches(0.35), font_size=13, bold=True, color=(0x1C, 0x1C, 0x1C))
            y_offset += 0.35
            
            # 每行3个指数卡片
            x_offset = 0.3
            for i, name in enumerate(names):
                if name not in indices:
                    continue
                info = indices[name]
                change = info['change']
                price = info['price']
                color_hex = (0x00, 0xAA, 0x66) if change >= 0 else (0xFF, 0x44, 0x44)
                
                # 卡片
                card = slide.shapes.add_shape(1, Inches(x_offset), Inches(y_offset), Inches(4.0), Inches(0.9))
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
                card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
                
                # 指数名称
                name_box = slide.shapes.add_textbox(Inches(x_offset) + Inches(0.1), Inches(y_offset) + Inches(0.05),
                                                   Inches(3.8), Inches(0.35))
                tf = name_box.text_frame
                p = tf.paragraphs[0]
                p.text = name
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
                
                # 价格
                price_box = slide.shapes.add_textbox(Inches(x_offset) + Inches(0.1), Inches(y_offset) + Inches(0.38),
                                                    Inches(2.5), Inches(0.35))
                tf2 = price_box.text_frame
                p2 = tf2.paragraphs[0]
                p2.text = f'{price:,.2f}'
                p2.font.size = Pt(14)
                p2.font.bold = True
                p2.font.color.rgb = RGBColor(0x1C, 0x1C, 0x1C)
                
                # 涨跌
                change_box = slide.shapes.add_textbox(Inches(x_offset) + Inches(2.5), Inches(y_offset) + Inches(0.38),
                                                      Inches(1.4), Inches(0.35))
                tf3 = change_box.text_frame
                p3 = tf3.paragraphs[0]
                emoji = '▲' if change >= 0 else '▼'
                p3.text = f'{emoji} {change:+.2f}%'
                p3.font.size = Pt(11)
                p3.font.bold = True
                p3.font.color.rgb = RGBColor(*color_hex)
                p3.alignment = PP_ALIGN.RIGHT
                
                x_offset += 4.2
                if x_offset > 10:
                    x_offset = 0.3
                    y_offset += 0.95
            y_offset += 1.0
    
    def _build_btc_technical_slide(self, prs, data):
        """BTC技术分析幻灯片"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._add_slide_header(slide, prs, '📈 BTC 技术分析',
                              f'趋势: {data["trend"]}  |  MA5: ${_fmt_val(data["ma5"])}  MA10: ${_fmt_val(data["ma10"])}  MA20: ${_fmt_val(data["ma20"])}')
        
        # 左：价格走势图
        chart_path = ChartGenerator.plot_btc_price_chart()
        if chart_path:
            self._add_image_to_slide(slide, chart_path, Inches(0.2), Inches(1.0), Inches(7.5), Inches(4.0))
        
        # 右：指标面板
        x_right = Inches(7.9)
        
        # 指标标题
        self._add_text_box(slide, '  技术指标', x_right, Inches(1.0),
                          Inches(4.8), Inches(0.35), font_size=13, bold=True, color=(0x1C, 0x1C, 0x1C))
        
        # RSI
        rsi = data['rsi']
        rsi_color = (0xFF, 0x44, 0x44) if (rsi and rsi > 70) else (0x00, 0xAA, 0x66) if (rsi and rsi < 30) else (0xFF, 0xCC, 0x00)
        rsi_box = slide.shapes.add_shape(1, x_right, Inches(1.4), Inches(4.8), Inches(0.7))
        rsi_box.fill.solid()
        rsi_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        rsi_box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        self._add_text_box(slide, f'RSI(14):  {f"{rsi:.1f}" if rsi else "N/A"}  —  {data["rsi_status"]}',
                          x_right + Inches(0.1), Inches(1.45), Inches(4.6), Inches(0.6),
                          font_size=11, color=rsi_color)
        
        # MACD
        macd_color = (0x00, 0xAA, 0x66) if data['histogram'] and data['histogram'] > 0 else (0xFF, 0x44, 0x44)
        macd_box = slide.shapes.add_shape(1, x_right, Inches(2.15), Inches(4.8), Inches(0.7))
        macd_box.fill.solid()
        macd_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        macd_box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        macd_str = f'${data["macd"]:.0f}' if data['macd'] else 'N/A'
        self._add_text_box(slide, f'MACD: {macd_str}  —  {data["macd_status"]}',
                          x_right + Inches(0.1), Inches(2.2), Inches(4.6), Inches(0.6),
                          font_size=11, color=macd_color)
        
        # 布林带
        bb_box = slide.shapes.add_shape(1, x_right, Inches(2.9), Inches(4.8), Inches(0.9))
        bb_box.fill.solid()
        bb_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        bb_box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        bb_text = f'布林带: Upper ${_fmt_val(data["bb_upper"])}  Middle ${_fmt_val(data["bb_middle"])}  Lower ${_fmt_val(data["bb_lower"])}'
        self._add_text_box(slide, bb_text,
                          x_right + Inches(0.1), Inches(2.95), Inches(4.6), Inches(0.8),
                          font_size=10, color=(0x44, 0x44, 0x44))
        
        # 支撑阻力
        sr_title = slide.shapes.add_shape(1, x_right, Inches(3.9), Inches(2.3), Inches(1.3))
        sr_title.fill.solid()
        sr_title.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
        sr_title.line.color.rgb = RGBColor(0x00, 0xAA, 0x66)
        self._add_text_box(slide, '🟢 支撑位',
                          x_right + Inches(0.1), Inches(3.95), Inches(2.1), Inches(0.3),
                          font_size=10, bold=True, color=(0x00, 0x77, 0x44))
        self._add_text_box(slide, f'${_fmt_val(data["recent_low"])} / ${_fmt_val(data["pivot"])}',
                          x_right + Inches(0.1), Inches(4.25), Inches(2.1), Inches(0.4),
                          font_size=10, color=(0x33, 0x33, 0x33))
        
        res_box = slide.shapes.add_shape(1, x_right + Inches(2.4), Inches(3.9), Inches(2.4), Inches(1.3))
        res_box.fill.solid()
        res_box.fill.fore_color.rgb = RGBColor(0xFF, 0xEE, 0xEE)
        res_box.line.color.rgb = RGBColor(0xFF, 0x44, 0x44)
        self._add_text_box(slide, '🔴 阻力位',
                          x_right + Inches(2.5), Inches(3.95), Inches(2.2), Inches(0.3),
                          font_size=10, bold=True, color=(0xCC, 0x22, 0x22))
        self._add_text_box(slide, f'${_fmt_val(data["recent_high"])} / $80,000',
                          x_right + Inches(2.5), Inches(4.25), Inches(2.2), Inches(0.4),
                          font_size=10, color=(0x33, 0x33, 0x33))
        
        # 信号面板
        signal_box = slide.shapes.add_shape(1, x_right, Inches(5.3), Inches(4.8), Inches(0.75))
        signal_box.fill.solid()
        signal_box.fill.fore_color.rgb = RGBColor(0x1C, 0x1C, 0x1C)
        signal_box.line.fill.background()
        sig_color = (0x00, 0xAA, 0x66) if data['overall'] == '偏多' else (0xFF, 0x44, 0x44) if data['overall'] == '偏空' else (0xFF, 0xCC, 0x00)
        self._add_text_box(slide, f'信号: {data["overall"]}  (看多{data["bull"]}/3  看空{data["bear"]}/3)',
                          x_right + Inches(0.1), Inches(5.35), Inches(4.6), Inches(0.6),
                          font_size=12, bold=True, color=sig_color)
    
    def _build_news_slide(self, prs, data):
        """市场新闻幻灯片"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._add_slide_header(slide, prs, '📰 今日市场新闻', '加密货币 · 宏观经济 · 地缘风险')
        
        news = data['news']
        if not news:
            self._add_text_box(slide, '暂无新闻数据', Inches(1), Inches(2), Inches(10), Inches(1),
                             font_size=14, color=(0x88, 0x88, 0x88), align='center')
            return
        
        y_offset = 1.1
        for i, n in enumerate(news):
            # 新闻卡片
            card = slide.shapes.add_shape(1, Inches(0.3), Inches(y_offset), Inches(12.67), Inches(1.1))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
            card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            
            # 序号标签
            num_tag = slide.shapes.add_shape(1, Inches(0.3), Inches(y_offset), Inches(0.4), Inches(1.1))
            num_tag.fill.solid()
            num_tag.fill.fore_color.rgb = RGBColor(0xF7, 0x93, 0x1A)
            num_tag.line.fill.background()
            self._add_text_box(slide, str(i+1), Inches(0.3), Inches(y_offset) + Inches(0.3),
                             Inches(0.4), Inches(0.5), font_size=14, bold=True,
                             color=(0xFF, 0xFF, 0xFF), align='center')
            
            # 标题
            self._add_text_box(slide, n['title'],
                              Inches(0.8), Inches(y_offset) + Inches(0.08), Inches(11.5), Inches(0.4),
                              font_size=12, bold=True, color=(0x1C, 0x1C, 0x1C))
            
            # 摘要
            summary = n['summary'][:150] + '...' if n['summary'] else ''
            self._add_text_box(slide, summary,
                              Inches(0.8), Inches(y_offset) + Inches(0.48), Inches(8), Inches(0.4),
                              font_size=9, color=(0x66, 0x66, 0x66))
            
            # 来源
            pub = n['pubDate'][:10] if n['pubDate'] else ''
            self._add_text_box(slide, f"来源: {n['provider']}  ·  {pub}",
                              Inches(9.0), Inches(y_offset) + Inches(0.48), Inches(3.8), Inches(0.4),
                              font_size=9, color=(0x99, 0x99, 0x99), align='right')
            
            y_offset += 1.2
    
    def _build_summary_slide(self, prs, data):
        """总结策略幻灯片"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._add_slide_header(slide, prs, '💼 交易策略 & 风险提示', f'综合信号: {data["overall"]}')
        
        # 组合对比图
        chart_path = ChartGenerator.plot_portfolio_bar(data['crypto'], data['metals'])
        if chart_path:
            self._add_image_to_slide(slide, chart_path, Inches(0.2), Inches(1.0), Inches(7.5), Inches(2.8))
        
        # 右侧：策略建议
        x_right = Inches(7.9)
        
        # 策略面板
        sig_color = (0x00, 0x77, 0x44) if data['overall'] == '偏多' else (0xCC, 0x22, 0x22) if data['overall'] == '偏空' else (0xCC, 0x88, 0x00)
        strat_box = slide.shapes.add_shape(1, x_right, Inches(1.0), Inches(4.8), Inches(2.8))
        strat_box.fill.solid()
        strat_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        strat_box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        
        self._add_text_box(slide, f'  信号: {data["overall"]}',
                          x_right, Inches(1.05), Inches(4.6), Inches(0.4),
                          font_size=14, bold=True, color=sig_color)
        
        current = data['current']
        recent_low = data['recent_low']
        recent_high = data['recent_high']
        
        self._add_text_box(slide,
                          f'  进场区域: ${current:,.0f}\n'
                          f'  止损位:   ${recent_low:,.0f}\n'
                          f'  目标位:   ${recent_high:,.0f}\n'
                          f'  仓位建议: 轻仓，不超过20%',
                          x_right, Inches(1.5), Inches(4.6), Inches(1.8),
                          font_size=11, color=(0x33, 0x33, 0x33))
        
        # 风险提示
        warn_box = slide.shapes.add_shape(1, Inches(0.3), Inches(4.0), Inches(12.67), Inches(1.3))
        warn_box.fill.solid()
        warn_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
        warn_box.line.color.rgb = RGBColor(0xFF, 0xBB, 0x33)
        
        self._add_text_box(slide, '⚠️ 风险提示',
                          Inches(0.4), Inches(4.1), Inches(12), Inches(0.35),
                          font_size=12, bold=True, color=(0xCC, 0x66, 0x00))
        self._add_text_box(slide,
                          '• 严格止损，不要扛单\n'
                          '• 控制仓位，杠杆不超过5x\n'
                          '• 本报告仅供参考，不构成投资建议',
                          Inches(0.4), Inches(4.45), Inches(12), Inches(0.8),
                          font_size=10, color=(0x66, 0x44, 0x00))
        
        # 底部免责声明
        self._add_text_box(slide,
                          '⚠️ 本报告由 Quinn 投资研究员自动生成，数据来源于 CoinGecko/yFinance/Alternative.me 等公开接口。'
                          '报告内容仅供参考，不构成任何投资建议。投资有风险，入市需谨慎。',
                          Inches(0.3), Inches(5.4), Inches(12.67), Inches(0.5),
                          font_size=8, color=(0xAA, 0xAA, 0xAA), align='center')
    
    def generate(self, output_path: str = None, external_data: dict = None) -> str:
        """
        生成完整PPT报告。
        external_data: dict 可选，传入则直接使用（避免 PPTX 内部重复下载触发 yfinance rate limit）
        """
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.oxml.ns import qn

        print('正在获取数据...')
        if external_data is not None:
            data = external_data
            print('使用外部传入数据，跳过内部下载和图表生成（图表已在 daily_push 中预生成）')
        else:
            data = self._get_data()
            print('正在生成图表...')
            # 预生成图表（避免并发问题）
            ChartGenerator.plot_btc_price_chart()
            ChartGenerator.plot_market_overview(data['crypto'], data['metals'], data['indices'])
            ChartGenerator.plot_portfolio_bar(data['crypto'], data['metals'])
        
        print('正在生成PPT...')
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        self._build_cover_slide(prs, data)
        self._build_market_overview_slide(prs, data)
        self._build_global_indices_slide(prs, data)
        self._build_btc_technical_slide(prs, data)
        self._build_news_slide(prs, data)
        self._build_summary_slide(prs, data)
        
        output_path = output_path or os.path.join(OUTPUT_DIR, f'quinn_report_{self.today.replace("-","")}.pptx')
        prs.save(output_path)
        print(f'报告已保存: {output_path}')
        return output_path


if __name__ == '__main__':
    import time
    t0 = time.time()
    generator = QuinnReportPPTX()
    path = generator.generate()
    print(f'总耗时: {time.time()-t0:.1f}s')
    print(f'输出: {path}')
